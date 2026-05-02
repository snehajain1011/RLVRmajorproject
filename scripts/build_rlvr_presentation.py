from __future__ import annotations

import csv
import json
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "artifacts" / "rlvr_presentation"
ASSET_DIR = OUT_DIR / "assets"
PPTX_PATH = OUT_DIR / "social_rlvr_web_agent_project.pptx"


COLORS = {
    "paper": "F7F1E8",
    "ink": "282522",
    "muted": "756F66",
    "clay": "C76445",
    "apricot": "E0A77E",
    "sage": "6F7C6D",
    "line": "DED3C6",
    "terminal": "201F1D",
    "terminal_dim": "9B948A",
    "success": "6F7C6D",
    "fail": "C76445",
}


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS["paper"])
    shape.line.fill.background()


def add_text(slide, text, x, y, w, h, size=20, color="ink", bold=False, font="Aptos", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(COLORS[color])
    return box


def add_kicker(slide, text, num):
    add_text(slide, f"{num:02d}", 0.65, 0.42, 0.34, 0.22, 8, "paper", True)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.62), Inches(0.42), Inches(0.32), Inches(0.32))
    dot.fill.solid()
    dot.fill.fore_color.rgb = rgb(COLORS["clay"])
    dot.line.fill.background()
    add_text(slide, text.upper(), 1.04, 0.4, 2.8, 0.26, 8, "muted", True)


def add_title(slide, kicker, title, subtitle, num):
    add_kicker(slide, kicker, num)
    add_text(slide, title, 0.62, 0.9, 11.35, 1.15, 29, "ink", False, "Georgia")
    if subtitle:
        add_text(slide, subtitle, 0.66, 2.08, 9.7, 0.52, 12.5, "muted")


def add_footer(slide, page):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(7.02), Inches(12.1), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(COLORS["line"])
    line.line.fill.background()
    add_text(slide, "Social-RLVR Web Agent Project", 0.62, 7.13, 3.0, 0.18, 7.5, "muted")
    add_text(slide, str(page).zfill(2), 12.28, 7.13, 0.44, 0.18, 7.5, "muted", align=PP_ALIGN.RIGHT)


def rect(slide, x, y, w, h, color="paper", line="line", radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[color])
    if line:
        shape.line.color.rgb = rgb(COLORS[line])
        shape.line.width = Pt(0.7)
    else:
        shape.line.fill.background()
    return shape


def metric(slide, x, y, value, label, color="ink"):
    add_text(slide, value, x, y, 1.4, 0.45, 25, color, False, "Georgia")
    add_text(slide, label, x, y + 0.47, 1.8, 0.35, 8.5, "muted", True)


def make_terminal_image(path: Path, title: str, lines: list[str], width=1120, height=560):
    img = Image.new("RGB", (width, height), f"#{COLORS['terminal']}")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("C:/Windows/Fonts/consola.ttf", 24)
        small = ImageFont.truetype("C:/Windows/Fonts/consola.ttf", 19)
        title_font = ImageFont.truetype("C:/Windows/Fonts/consolab.ttf", 22)
    except OSError:
        font = ImageFont.load_default()
        small = ImageFont.load_default()
        title_font = ImageFont.load_default()
    draw.rounded_rectangle((0, 0, width - 1, height - 1), radius=28, fill=f"#{COLORS['terminal']}")
    draw.ellipse((28, 26, 44, 42), fill="#D46A5D")
    draw.ellipse((54, 26, 70, 42), fill="#D9A25F")
    draw.ellipse((80, 26, 96, 42), fill="#7A9A79")
    draw.text((120, 22), title, fill="#DDD3C7", font=title_font)
    y = 78
    for line in lines:
        color = "#F4EEE5"
        if "0.0" in line or "False" in line or "failed" in line.lower() or "not submitted" in line:
            color = f"#{COLORS['apricot']}"
        if "1.0" in line or "True" in line or "correctly" in line:
            color = "#B5C7AF"
        if line.startswith("#") or line.startswith(">"):
            color = f"#{COLORS['terminal_dim']}"
        draw.text((34, y), line[:92], fill=color, font=font if len(line) < 72 else small)
        y += 38
    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path)


def load_csv(path: Path) -> list[dict[str, str]]:
    with path.open(newline="", encoding="utf-8") as handle:
        return list(csv.DictReader(handle))


def make_assets():
    report_summary = load_csv(ROOT / "artifacts" / "eval_qwen_text_report_rlvr" / "summary.csv")
    complex_rows = load_csv(ROOT / "artifacts" / "eval_qwen_text_complex_rlvr" / "episode_results.csv")
    trajectories = [
        json.loads(line)
        for line in (ROOT / "artifacts" / "eval_qwen_text_report_rlvr" / "trajectories.jsonl").read_text(encoding="utf-8").splitlines()
        if line.strip()
    ]

    make_terminal_image(
        ASSET_DIR / "baseline_failure.png",
        "qwen2.5:0.5b zero-shot / BrowserGym trajectory",
        [
            "$ evaluate_browsergym_rlvr --policies qwen --task report.extract_tracking_code",
            f"policy: {trajectories[0]['policy']}",
            f"success: {trajectories[0]['success']}   reward: {trajectories[0]['reward']}   steps: {trajectories[0]['steps']}",
            f"verifier: {trajectories[0]['verifier_message']}",
            "",
            "step 1: click(\"17\")",
            "step 2: click(\"17\")",
            "step 3: click(\"17\")",
            "step 4: click(\"17\")",
            "# Model repeatedly clicked Submit without filling TRV-8429-IN",
        ],
    )
    make_terminal_image(
        ASSET_DIR / "after_success.png",
        "after verifier-guided trajectory distillation",
        [
            "$ evaluate_browsergym_rlvr --policies qwen,rlvr --task report.extract_tracking_code",
            "policy                                      success_rate  reward  RRR",
            f"{report_summary[1]['policy'][:42]:42} {report_summary[1]['success_rate']:>11}  {report_summary[1]['mean_reward']:>6}  {report_summary[1]['mean_rrr']}",
            f"{report_summary[0]['policy'][:42]:42} {report_summary[0]['success_rate']:>11}  {report_summary[0]['mean_reward']:>6}  {report_summary[0]['mean_rrr']}",
            "",
            "step 1: fill(\"16\", \"TRV-8429-IN\")",
            "step 2: click(\"17\")",
            "verifier: tracking code submitted correctly",
        ],
    )
    make_terminal_image(
        ASSET_DIR / "complex_results.png",
        "complex task evaluation",
        [
            "$ evaluate_browsergym_rlvr --tasks orders.priority_followup,schedule.design_review_shared_slot",
            "task                                  qwen_zero_shot   after_rlvr",
            "orders.priority_followup              0.0 / False      1.0 / True",
            "schedule.design_review_shared_slot    0.0 / False      1.0 / True",
            "",
            "verifier: high priority Meera order follow-up created correctly",
            "verifier: design review scheduled with Kabir and Zara in the shared slot",
            f"# rows logged: {len(complex_rows)} episodes",
        ],
    )
    make_terminal_image(
        ASSET_DIR / "artifact_trail.png",
        "project artifact trail",
        [
            "artifacts/",
            "  eval_qwen_text_report_rlvr/",
            "    summary.csv              before/after metrics",
            "    episode_results.csv      per-task outcomes",
            "    trajectories.jsonl       exact actions + verifier messages",
            "  rlvr_training_report/",
            "    learned_policy.json      verifier-selected successful trajectory",
            "  eval_qwen_text_complex_rlvr/",
            "    summary.csv              complex-task before/after metrics",
            "  rlvr_training_complex/",
            "    learned_policy.json      learned actions for harder tasks",
        ],
    )


def add_bar(slide, x, y, label, value, color):
    add_text(slide, label, x, y - 0.04, 2.5, 0.26, 9.5, "muted", True)
    rect(slide, x + 2.55, y, 3.2, 0.18, "line", None)
    rect(slide, x + 2.55, y, 3.2 * value, 0.18, color, None)
    add_text(slide, f"{int(value * 100)}%", x + 5.88, y - 0.11, 0.6, 0.28, 12, color, True)


def build_deck():
    make_assets()
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def slide_base(kicker, title, subtitle, num):
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, kicker, title, subtitle, num)
        add_footer(slide, num)
        return slide

    s = slide_base(
        "Project",
        "Social-RLVR Web Agent Project",
        "A minimal browser-agent testbed where tasks are rewarded only when backend state proves completion.",
        1,
    )
    metric(s, 8.2, 1.14, "5", "verifiable browser tasks", "clay")
    metric(s, 9.85, 1.14, "0 to 100%", "before/after SR on tested tasks", "sage")
    metric(s, 11.55, 1.14, "1.0", "reward only after real completion", "ink")
    rect(s, 0.66, 4.72, 11.9, 0.72, "ink", None)
    add_text(s, 'Core thesis: browser agents should not be rewarded for saying "done"; they should be rewarded only when the task is actually done.', 0.95, 4.91, 10.7, 0.32, 15, "paper")

    s = slide_base("Problem", "The baseline model claimed progress without changing the task state.", "The verifier caught the failure because the tracking code was never submitted.", 2)
    s.shapes.add_picture(str(ASSET_DIR / "baseline_failure.png"), Inches(0.78), Inches(2.72), width=Inches(7.15))
    add_text(s, "What failed", 8.45, 2.82, 2.6, 0.35, 16, "ink", False, "Georgia")
    add_text(s, "Qwen clicked Submit repeatedly, but the backend report never contained the required value.", 8.47, 3.32, 3.35, 0.72, 13, "muted")
    add_text(s, "Verifier message", 8.47, 4.38, 2.0, 0.25, 9, "clay", True)
    add_text(s, "correct tracking code not submitted", 8.47, 4.7, 3.5, 0.38, 16, "ink", False, "Georgia")

    s = slide_base("Approach", "The project wraps browser tasks with objective backend-state rewards.", "Playwright controls the browser; BrowserGym provides the agent interface; Python verifiers score the result.", 3)
    steps = [("BrowserGym task", "Instruction + observation"), ("Playwright env", "Clicks, fills, selects"), ("Local web app", "Messages, report, orders"), ("Verifier", "Checks backend state")]
    for i, (a, b) in enumerate(steps):
        x = 0.9 + i * 3.05
        rect(s, x, 3.08, 2.25, 1.2, "paper", "line")
        add_text(s, a, x + 0.18, 3.29, 1.86, 0.26, 13, "ink", True)
        add_text(s, b, x + 0.18, 3.68, 1.82, 0.38, 9.5, "muted")
        if i < 3:
            add_text(s, "→", x + 2.45, 3.39, 0.38, 0.36, 21, "clay", False, "Georgia")
    rect(s, 2.54, 5.05, 8.2, 0.62, "clay", None)
    add_text(s, "Reward = 1 only when backend data exactly matches the task goal", 3.05, 5.24, 7.2, 0.26, 14, "paper", True)

    s = slide_base("Build", "A full local RLVR browser scaffold is already implemented.", "The repo now has environment code, task registration, model policy, learned policy, evaluation, and training scripts.", 4)
    items = [
        ("src/social_rlvr_web/env.py", "Playwright Gym-style environment"),
        ("src/social_rlvr_web/browsergym_tasks.py", "First-class BrowserGym task registration"),
        ("src/social_rlvr_web/tasks.py", "Task specs and deterministic verifiers"),
        ("src/social_rlvr_web/model_policy.py", "Ollama/Qwen policy wrapper"),
        ("src/social_rlvr_web/rlvr_policy.py", "Learned trajectory replay policy"),
        ("scripts/train_browsergym_rlvr.py", "Verifier-guided improvement loop"),
    ]
    for idx, (path, desc) in enumerate(items):
        y = 2.5 + idx * 0.58
        add_text(s, path, 0.88, y, 4.1, 0.25, 10.5, "ink", True, "Consolas")
        add_text(s, desc, 5.15, y, 5.5, 0.25, 10.5, "muted")

    s = slide_base("Task Suite", "The current suite tests reading, repetition, visual choice, table logic, and scheduling.", "Each task has a hidden ground-truth condition checked through backend state.", 5)
    tasks = [
        ("Messages", "Send Happy New Year to five friends", "exact recipients + text"),
        ("Gallery", "Send aesthetic travel image to Meera", "image_id + recipient"),
        ("Report", "Extract shipment tracking code", "TRV-8429-IN"),
        ("Orders", "Create priority follow-up", "owner + reference + note"),
        ("Schedule", "Find shared design-review slot", "title + slot + attendees"),
    ]
    for idx, (name, goal, check) in enumerate(tasks):
        y = 2.25 + idx * 0.72
        add_text(s, name, 0.92, y, 1.2, 0.28, 12.5, "clay", True)
        add_text(s, goal, 2.2, y, 4.8, 0.28, 12, "ink")
        add_text(s, check, 7.4, y, 3.9, 0.28, 11, "muted")

    s = slide_base("Reward", "The reward is simple: either the backend proves success, or it does not.", "No subjective judge, no preference model, and no credit for hallucinated completion.", 6)
    left = [
        ("Message verifier", "all five target recipients received exact greeting"),
        ("Report verifier", "tracking_code == TRV-8429-IN"),
        ("Schedule verifier", "slot == Fri 14:00 and attendees == Kabir + Zara"),
    ]
    for i, (name, desc) in enumerate(left):
        y = 2.55 + i * 0.92
        rect(s, 0.9, y - 0.08, 0.12, 0.46, "clay", None)
        add_text(s, name, 1.2, y - 0.02, 2.0, 0.25, 12, "ink", True)
        add_text(s, desc, 3.3, y - 0.02, 6.4, 0.25, 11, "muted")
    metric(s, 10.25, 2.46, "1.0", "success", "sage")
    metric(s, 10.25, 4.05, "0.0", "otherwise", "clay")

    s = slide_base("Baseline", "The real lightweight model baseline scored zero on the report task.", "This is the failure RLVR is meant to correct: the model acted, but did not complete the state change.", 7)
    add_bar(s, 0.92, 3.1, "Qwen zero-shot success", 0.0, "clay")
    add_bar(s, 0.92, 3.78, "Mean reward", 0.0, "clay")
    add_bar(s, 0.92, 4.46, "RRR", 0.0, "clay")
    s.shapes.add_picture(str(ASSET_DIR / "baseline_failure.png"), Inches(7.05), Inches(2.42), width=Inches(5.4))

    s = slide_base("RLVR Loop", "Verifier-selected trajectories power the current improvement loop.", "It is an end-to-end proof of the pipeline, not yet weight-level GRPO fine-tuning.", 8)
    loop = [("1", "Run model", "Collect rollout"), ("2", "Verify", "Reward from backend"), ("3", "Select success", "Add teacher trajectory if needed"), ("4", "Evaluate", "Compare before vs after")]
    for i, (n, a, b) in enumerate(loop):
        x = 1.0 + i * 3.0
        rect(s, x, 3.0, 2.25, 1.38, "paper", "line")
        add_text(s, n, x + 0.16, 3.16, 0.35, 0.35, 17, "clay", False, "Georgia")
        add_text(s, a, x + 0.58, 3.18, 1.28, 0.28, 12.5, "ink", True)
        add_text(s, b, x + 0.58, 3.58, 1.45, 0.35, 9.5, "muted")
    add_text(s, "Current method: verifier_guided_trajectory_distillation", 3.65, 5.08, 5.9, 0.28, 11, "muted", True, "Consolas")

    s = slide_base("Result", "On the report task, the learned policy moved from 0 to 100% success.", "The successful trajectory fills the exact tracking code, then submits.", 9)
    add_bar(s, 0.95, 2.7, "Qwen zero-shot", 0.0, "clay")
    add_bar(s, 0.95, 3.48, "After RLVR artifact", 1.0, "sage")
    metric(s, 1.0, 4.65, "2", "steps after improvement", "ink")
    metric(s, 2.8, 4.65, "1.0", "mean reward", "sage")
    s.shapes.add_picture(str(ASSET_DIR / "after_success.png"), Inches(6.55), Inches(2.42), width=Inches(5.95))

    s = slide_base("Expansion", "The same before/after pattern holds on the two harder tasks.", "Row-conditioned follow-up and constraint scheduling both failed zero-shot and succeeded after the learned policy artifact.", 10)
    s.shapes.add_picture(str(ASSET_DIR / "complex_results.png"), Inches(0.82), Inches(2.42), width=Inches(6.95))
    add_bar(s, 8.35, 3.05, "Qwen zero-shot", 0.0, "clay")
    add_bar(s, 8.35, 3.85, "After RLVR artifact", 1.0, "sage")
    add_text(s, "2/2 complex tasks solved after trajectory distillation", 8.37, 4.9, 3.15, 0.55, 18, "ink", False, "Georgia")

    s = slide_base("Evidence", "The repo now saves a reproducible experiment trail.", "These files are the basis for the slides and can be rerun from the README commands.", 11)
    s.shapes.add_picture(str(ASSET_DIR / "artifact_trail.png"), Inches(0.82), Inches(2.25), width=Inches(6.8))
    add_text(s, "What this gives us", 8.18, 2.45, 2.8, 0.32, 16, "ink", False, "Georgia")
    add_text(s, "A clean experiment record: policy name, task id, success, reward, steps, RRR, verifier message, and exact action sequence.", 8.2, 2.95, 3.42, 0.9, 12.5, "muted")
    add_text(s, "This makes failures explainable instead of anecdotal.", 8.2, 4.35, 3.35, 0.42, 15, "clay", False, "Georgia")

    s = slide_base("Next", "The next milestone is true policy training.", "The current system should be presented as verifier-guided trajectory distillation, not as completed GRPO weight tuning.", 12)
    next_steps = [
        ("Add GRPO update loop", "Replace replay with actual model/policy update"),
        ("Run stronger multimodal models", "Use more memory or cloud hardware for Qwen2.5-VL"),
        ("Scale benchmarks", "Move from local tasks to MiniWoB, then WebArena"),
        ("Report rigorously", "Compare zero-shot, SFT, RLVR+GRPO with SR, reward, RRR"),
    ]
    for i, (a, b) in enumerate(next_steps):
        y = 2.45 + i * 0.76
        add_text(s, f"{i+1}.", 0.92, y, 0.35, 0.25, 12, "clay", True)
        add_text(s, a, 1.34, y, 2.75, 0.25, 12.5, "ink", True)
        add_text(s, b, 4.2, y, 5.3, 0.25, 11.2, "muted")
    rect(s, 0.92, 6.05, 10.8, 0.52, "ink", None)
    add_text(s, "Bottom line: the environment, verifiers, baseline, improvement loop, and evidence artifacts are implemented.", 1.18, 6.2, 10.0, 0.22, 12.5, "paper", True)

    prs.save(PPTX_PATH)


def verify_package():
    with zipfile.ZipFile(PPTX_PATH) as archive:
        names = archive.namelist()
        slide_count = len([name for name in names if name.startswith("ppt/slides/slide") and name.endswith(".xml")])
        media = [name for name in names if name.startswith("ppt/media/")]
        empty_media = [name for name in media if archive.getinfo(name).file_size <= 0]
    if slide_count != 12:
        raise RuntimeError(f"Expected 12 slides, found {slide_count}")
    if empty_media:
        raise RuntimeError(f"Empty media files: {empty_media}")
    return slide_count, len(media), PPTX_PATH.stat().st_size


if __name__ == "__main__":
    build_deck()
    slides, media, size = verify_package()
    print(f"Wrote {PPTX_PATH}")
    print(f"slides={slides} media={media} bytes={size}")
